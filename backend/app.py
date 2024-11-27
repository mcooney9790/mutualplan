from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pandas as pd
import calendar
from pptx.enum.text import PP_ALIGN
import os

app = Flask(__name__)
CORS(app)

# Function to create calendar slide for each month
def create_calendar_slide(prs, year, month, events):
    # Add a new slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide

    # Add title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title.text_frame.text = f"{calendar.month_name[month]} {year}"
    title.text_frame.paragraphs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Create a table for the calendar grid
    rows, cols = 7, 7
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    for i in range(cols):
        table.columns[i].width = Inches(1.25)

    # Add day names to the first row
    day_names = ["Sun", "Mon", "Tue", "Wed", "Thu", "Fri", "Sat"]
    for col_idx, day_name in enumerate(day_names):
        cell = table.cell(0, col_idx)
        cell.text = day_name
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.text_frame.paragraphs[0].font.bold = True

    # Populate the table with days and events
    month_calendar = calendar.monthcalendar(year, month)
    for row_idx, week in enumerate(month_calendar):
        for col_idx, day in enumerate(week):
            cell = table.cell(row_idx + 1, col_idx)
            if day != 0:
                cell.text = str(day)
                for _, event in events[events['Date'].dt.day == day].iterrows():
                    cell.text += f"\n{event['nextSteps']}"

# Endpoint to generate and download PowerPoint
@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    data = request.json
    timeline = data.get('timeline', [])

    df = pd.DataFrame(timeline)
    df['Date'] = pd.to_datetime(df['date'])

    prs = Presentation()

    df['Year'] = df['Date'].dt.year
    df['Month'] = df['Date'].dt.month
    grouped_data = df.groupby(['Year', 'Month'])

    for (year, month), events in grouped_data:
        create_calendar_slide(prs, year, month, events)

    output_path = "timeline_calendar.pptx"
    prs.save(output_path)

    # Send the file as a downloadable response
    return send_file(output_path, as_attachment=True, download_name="timeline_calendar.pptx")

if __name__ == '__main__':
    app.run(debug=True)
